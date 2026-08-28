from pptx import Presentation


def processar_pptx(caminho: str):
    apresentacao = Presentation(caminho)
    blocos = []
    for i, slide in enumerate(apresentacao.slides, start=1):
        textos = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if textos:
            blocos.append({"slide": i, "texto": "\n".join(textos)})
    return blocos
