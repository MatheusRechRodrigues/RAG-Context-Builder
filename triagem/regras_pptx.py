from pptx import Presentation


def triar_pptx(doc):
    try:
        apresentacao = Presentation(doc.caminho)
    except Exception as e:
        doc.nivel = 3
        doc.motivos.append(f"Falha ao abrir PPTX (corrompido ou protegido): {e}")
        return doc

    slides = list(apresentacao.slides)
    total_slides = len(slides)

    if total_slides == 0:
        doc.nivel = 3
        doc.motivos.append("Apresentacao sem slides")
        return doc

    slides_com_texto = 0
    for slide in slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slides_com_texto += 1
                break

    percentual = slides_com_texto / total_slides

    if percentual < 0.3:
        doc.nivel = 3
        doc.pipeline = "ocr"
        doc.motivos.append(f"Apenas {percentual:.0%} dos slides tem texto extraivel")
        return doc

    if percentual < 0.7:
        doc.nivel = 2
        doc.pipeline = "pptx_texto"
        doc.motivos.append(f"{percentual:.0%} dos slides tem texto extraivel")
        return doc

    doc.nivel = 1
    doc.pipeline = "pptx_texto"
    return doc
